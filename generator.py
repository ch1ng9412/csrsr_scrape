import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import glob

def iter_text_shapes(shapes):
    """
    遞迴尋找所有包含文字的物件。
    (解決文字框被「群組 Group」或放在「表格 Table」中導致抓不到的問題)
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_text_shapes(shape.shapes)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell
        else:
            yield shape

def iter_picture_shapes(shapes):
    """
    【全新升級】遞迴尋找所有圖片物件。
    (解決圖片被「群組 Group」在一起時，程式抓不到佔位圖的問題)
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_picture_shapes(shape.shapes)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape

def replace_text_in_slide(slide, replace_dict):
    """替換投影片中的文字，使用「段落級別」與「群組遞迴」替換"""
    for shape in iter_text_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = paragraph.text
            needs_replace = False
            
            for key in replace_dict.keys():
                if key in full_text:
                    needs_replace = True
                    break
            
            if needs_replace:
                for key, val in replace_dict.items():
                    full_text = full_text.replace(key, str(val))
                
                # 保留原本的字體格式
                if len(paragraph.runs) > 0:
                    paragraph.runs[0].text = full_text
                    for i in range(1, len(paragraph.runs)):
                        paragraph.runs[i].text = ""

def get_image_path(img_dir, slope_id, suffix):
    """根據指定的後綴尋找對應的圖片"""
    pattern = os.path.join(img_dir, f"{slope_id}{suffix}.*")
    files = glob.glob(pattern)
    valid_exts = ('.jpg', '.jpeg', '.png')
    
    for f in files:
        if f.lower().endswith(valid_exts):
            return f
    return None

def replace_images_in_slide(slide, new_images):
    """
    依照圖層順序替換圖片，並完美保留「原本佔位圖片的圖層位置(Z-order)」。
    支援處理被群組(Grouped)的圖片。
    """
    # 使用升級版的遞迴掃描，把投影片中所有圖片都抓出來
    pic_shapes = list(iter_picture_shapes(slide.shapes))
            
    for i, old_shape in enumerate(pic_shapes):
        if i < len(new_images):
            img_path = new_images[i]
            if img_path and os.path.exists(img_path):
                left, top = old_shape.left, old_shape.top
                width, height = old_shape.width, old_shape.height
                
                # 1. 插入新圖片 (預設加在投影片最上層)
                new_shape = slide.shapes.add_picture(img_path, left, top, width, height)
                
                # 2. 【關鍵黑科技】調整 Z-order 與 群組層級
                # 不管舊圖片是在群組裡還是獨立存在，直接將新圖片的 XML 塞到舊圖片的 XML 前面
                old_element = old_shape._element
                new_element = new_shape._element
                old_element.addprevious(new_element)
                
                # 3. 刪除舊圖片
                old_element.getparent().remove(old_element)
            else:
                print(f"  [警告] 找不到第 {i + 1} 張佔位圖對應的圖片，跳過替換。")

def main():
    # ⚠️ 貼心提醒：請確保您的新模板檔案名稱改為 'Template.pptx'，或在此處修改檔名
    template_path = 'Template.pptx' 
    excel_path = 'Data.xlsx'
    img_dir = 'img' 

    if not os.path.exists(template_path):
        print(f"錯誤：找不到模板檔案 '{template_path}'")
        print("💡 提示：請確認您的新模板已經重新命名為 'Template.pptx'")
        return
    if not os.path.exists(excel_path):
        print(f"錯誤：找不到資料表 '{excel_path}'")
        return

    print("讀取資料表中...")
    try:
        df = pd.read_excel(excel_path).fillna('')
    except Exception as e:
        print(f"讀取 Excel 失敗: {e}")
        return

    os.makedirs("Output", exist_ok=True)

    for index, row in df.iterrows():
        slope_id = str(row.get('SLOPE_ID', f'未命名_{index}'))
        print(f"正在處理: {slope_id} ...")
        
        prs = Presentation(template_path)
        slide = prs.slides[0]

        # 處理 LEVEL (位移速率分級) 的數值轉換
        raw_level = str(row.get('LEVEL', '')).strip()
        if raw_level.endswith('.0'):
            raw_level = raw_level[:-2]
        
        level_mapping = {'1': '小', '2': '中', '3': '大'}
        level_str = level_mapping.get(raw_level, raw_level)

        # 建立文字替換字典
        replace_dict = {
            "{DATE}": row.get('DATE', ''),
            "2018/7 - 2024/4": row.get('DATE', ''), # 防呆：處理新模板中不小心被寫死的日期
            "2018/7-2024/4": row.get('DATE', ''),   # 兼容沒有空格的版本
            "{SLOPE_ID}": row.get('SLOPE_ID', ''),
            "{BASIN}": row.get('BASIN', ''),
            "{CountyName}": row.get('CountyName', ''),
            "{TownName}": row.get('TownName', ''),
            "{VillName}": row.get('VillName', ''),
            "{LEVEL}": level_str,
            "{AEAR}": row.get('AREA', ''), 
            "{AREA}": row.get('AREA', ''), 
            "{TWD97_E}": row.get('TWD97_E', ''),
            "{TWD97_N}": row.get('TWD97_N', '')
        }

        # 1. 執行文字替換
        replace_text_in_slide(slide, replace_dict)

        # 2. 執行圖片替換邏輯
        # 💡 【重要注意】因為您重新排版了模板圖片，PPT 底層的圖片順序極可能已經改變！
        # 如果跑出來後發現這三張圖片互相跑錯了位置（例如分布圖跑到地形圖的框框），
        # 請根據實際錯置的狀況，直接「上下調換」下面這三行 -4, -1, -3 的順序即可！
        target_suffixes = [
            "-3", # 對應畫面底層第 1 張圖片 ⭢ 土地類別分布圖
            "-4", # 對應畫面底層第 2 張圖片 ⭢ 殘坡區域衛星影像
            "-1"  # 對應畫面底層第 3 張圖片 ⭢ 1/25000地形圖
        ]
        
        new_images = []
        for suffix in target_suffixes:
            img_path = get_image_path(img_dir, slope_id, suffix)
            new_images.append(img_path)
        
        # 替換圖片並保留原有圖層與群組狀態
        replace_images_in_slide(slide, new_images)

        # 3. 存檔輸出
        output_filename = os.path.join("Output", f"{slope_id}_殘坡基本資料.pptx")
        prs.save(output_filename)
        print(f"  ✅ 完成產出: {output_filename}\n")

    print("🎉 所有簡報生成完畢！請至 Output 資料夾查看。")

if __name__ == "__main__":
    main()